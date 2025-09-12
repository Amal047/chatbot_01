# utils.py
import os
import re
import time
import pickle
import logging
from typing import List, Optional, Iterable, Tuple, Dict, Any

from sentence_transformers import SentenceTransformer
from groq import Groq
from dotenv import load_dotenv
from pdf2image import convert_from_path
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
import numpy as np
import pytesseract

from spellchecker import SpellChecker
from langdetect import detect
from deep_translator import GoogleTranslator

from app import db
from app.models import ChatHistory



# Load env
load_dotenv()
logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

# Config
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GROQ_MODEL = os.getenv("GROQ_MODEL", "llama-3.1-8b-instant")
ENABLE_SPELLCHECK = os.getenv("ENABLE_SPELLCHECK", "false").lower() in ("1", "true", "yes")
ENABLE_TRANSLATE = os.getenv("ENABLE_TRANSLATE", "false").lower() in ("1", "true", "yes")
EMBEDDING_MODEL_NAME = os.getenv("EMBEDDING_MODEL", "all-mpnet-base-v2")
EMBED_CACHE_PATH = os.getenv("EMBED_CACHE_PATH", "embed_cache.pkl")
OCR_DPI = int(os.getenv("OCR_DPI", "200"))

# Initialize
spell = SpellChecker() if ENABLE_SPELLCHECK else None
model = SentenceTransformer(EMBEDDING_MODEL_NAME)
groq_client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None

# Basic utilities
def _safe_read_text_from_docx(path: str) -> str:
    doc = DocxDocument(path)
    texts = []
    for para in doc.paragraphs:
        texts.append(para.text)
    return "\n".join(texts)


def _safe_read_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def _safe_read_text_from_xlsx(path: str) -> str:
    df = pd.read_excel(path, dtype=str, engine="openpyxl")
    # join non-null cells per row
    rows = df.fillna("").astype(str).apply(lambda r: " | ".join(r.values.tolist()), axis=1).tolist()
    return "\n".join(rows)


def _safe_read_text_from_csv(path: str) -> str:
    df = pd.read_csv(path, dtype=str)
    rows = df.fillna("").astype(str).apply(lambda r: " | ".join(r.values.tolist()), axis=1).tolist()
    return "\n".join(rows)


def extract_text(file_path: str) -> str:
    """
    Extract text using appropriate handler based on file extension.
    PDF: PyPDF2 + OCR fallback.
    DOCX / PPTX / CSV / XLSX / TXT: appropriate handlers.
    """
    ext = file_path.rsplit(".", 1)[-1].lower()
    try:
        if ext == "pdf":
            # Use PyPDF2 first
            reader = PdfReader(file_path)
            txt = []
            for page in reader.pages:
                text = page.extract_text() or ""
                txt.append(text)
            content = "\n".join(txt).strip()
            if len(content) < 50:
                # OCR fallback
                images = convert_from_path(file_path, dpi=OCR_DPI)
                ocr_texts = [pytesseract.image_to_string(img) for img in images]
                content = "\n".join([content] + ocr_texts).strip()
            return content
        elif ext in ("docx",):
            return _safe_read_text_from_docx(file_path)
        elif ext in ("pptx",):
            return _safe_read_text_from_pptx(file_path)
        elif ext in ("xlsx", "xls"):
            return _safe_read_text_from_xlsx(file_path)
        elif ext in ("csv",):
            return _safe_read_text_from_csv(file_path)
        elif ext in ("txt",):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        else:
            # Fallback: try reading as text
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
    except Exception as e:
        logger.exception("extract_text failed for %s: %s", file_path, e)
        raise


# -------------------
# Text cleaning / normalization
# -------------------
def clean_whitespace(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def fix_spelling_if_needed(text: str) -> str:
    if not ENABLE_SPELLCHECK or not spell:
        return text
    words = text.split()
    corrected = []
    for w in words:
        # avoid correcting tokens like emails, urls, code-like tokens
        if "@" in w or "." in w or "/" in w or len(w) < 3:
            corrected.append(w)
            continue
        corrected.append(spell.correction(w) or w)
    return " ".join(corrected)


def normalize_language(text: str, target_lang="en") -> str:
    if not ENABLE_TRANSLATE:
        return text
    try:
        sample = text[:500].strip()
        if not sample:
            return text
        lang = detect(sample)
        if lang and lang != target_lang:
            logger.info("Translating detected language %s -> %s", lang, target_lang)
            return GoogleTranslator(source="auto", target=target_lang).translate(text)
    except Exception:
        logger.exception("Language detection/translation failed")
    return text


# -------------------
# Smart chunking
# -------------------
_SENTENCE_SPLIT_RE = re.compile(r"(?<=[.!?])\s+")


def _split_into_sentences(text: str) -> List[str]:
    return [s.strip() for s in _SENTENCE_SPLIT_RE.split(text) if s.strip()]


def smart_split_text(text: str, max_chars: int = 800, min_chars: int = 100) -> List[str]:
    """
    Split text into semantically coherent chunks roughly of size `max_chars`.
    - Splits into sentences then merges until the chunk length approaches max_chars.
    - Ensures no chunk below min_chars unless it's the last one.
    """
    text = clean_whitespace(text)
    sentences = _split_into_sentences(text)
    chunks = []
    current = []
    cur_len = 0
    for s in sentences:
        if cur_len + len(s) + 1 <= max_chars or not current:
            current.append(s)
            cur_len += len(s) + 1
        else:
            chunk_text = " ".join(current).strip()
            if len(chunk_text) >= min_chars:
                chunks.append(chunk_text)
            else:
                # if short, append to previous chunk if exists, else keep
                if chunks:
                    chunks[-1] += " " + chunk_text
                else:
                    chunks.append(chunk_text)
            current = [s]
            cur_len = len(s) + 1
    if current:
        chunks.append(" ".join(current).strip())
    return chunks


# -------------------
# Embeddings & caching
# -------------------
# load or create embed cache
try:
    _embed_cache = pickle.load(open(EMBED_CACHE_PATH, "rb"))
except Exception:
    _embed_cache = {}

def _embed_key(text: str) -> str:
    return str(hash(text))

def embed_text(text: str) -> np.ndarray:
    """
    Return a numpy float32 embedding for the text. Caches embeddings on disk.
    """
    key = _embed_key(text)
    if key in _embed_cache:
        vec = np.array(_embed_cache[key], dtype="float32")
        return vec
    vec = model.encode([text], convert_to_numpy=True)[0].astype("float32")
    _embed_cache[key] = vec.tolist()
    # persist lazily
    try:
        pickle.dump(_embed_cache, open(EMBED_CACHE_PATH, "wb"))
    except Exception:
        logger.exception("Failed to persist embed cache")
    return vec


def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    if a is None or b is None:
        return 0.0
    a_norm = np.linalg.norm(a)
    b_norm = np.linalg.norm(b)
    if a_norm == 0 or b_norm == 0:
        return 0.0
    return float(np.dot(a, b) / (a_norm * b_norm))


# -------------------
# LLM / Groq helpers
# -------------------
import backoff
from requests import HTTPError

def _groq_create_chat(prompt: str, max_tokens: int = 512, temperature: float = 0.2) -> str:
    if not groq_client:
        raise RuntimeError("Groq client not configured (GROQ_API_KEY missing).")
    # simple retry/backoff for transient failures
    @backoff.on_exception(backoff.expo, Exception, max_tries=4)
    def _call():
        resp = groq_client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=temperature,
            max_tokens=max_tokens,
        )
        return resp
    res = _call()
    try:
        return res.choices[0].message.content.strip()
    except Exception:
        logger.exception("Unexpected Groq response format")
        raise


def ask_groq(prompt: str, max_tokens: int = 512, temperature: float = 0.2) -> str:
    try:
        return _groq_create_chat(prompt, max_tokens=max_tokens, temperature=temperature)
    except Exception as e:
        logger.exception("Groq request failed")
        return ""


def ask_local_llm_with_context(
    query: str,
    doc_context_texts: Optional[Iterable[str]] = None,
    chat_context_texts: Optional[Iterable[str]] = None,
    top_k_docs: int = 5,
    doc_threshold: float = 0.3,
    chat_threshold: float = 0.6,
) -> Tuple[Optional[str], str]:
    """
    Build a strict prompt using docs + chat history and ask Groq (or local LLM).
    Returns (answer_text, source) where source indicates 'local' or 'groq' or 'none'
    """
    docs = list(doc_context_texts or [])
    chats = list(chat_context_texts or [])

    # Build prompt: include top relevant docs only (we already used vectorstore search so docs should be relevant)
    prompt_parts = []
    if docs:
        prompt_parts.append("Document Context:\n" + "\n\n".join(docs[:top_k_docs]))
    if chats:
        prompt_parts.append("Previous Chat Context:\n" + "\n\n".join(chats[:3]))

    prompt_parts.append(f"User Question:\n{query}\n\nStrict instructions:\n"
                        "- Only use the provided context. If the answer is not explicitly in the context, reply exactly: 'Not mentioned in the document.'\n"
                        "- Do NOT hallucinate. Keep answers concise and factual.\n"
                        "- If the answer is a list, provide a bullet list.\n"
                        "- If you must quote a chunk, show which chunk it came from.\n")

    full_prompt = "\n\n".join(prompt_parts)

    # Call Groq (or substitute local LLM if you integrate one later)
    answer = ask_groq(full_prompt)
    source = "groq" if answer else "none"
    return answer, source


# -------------------
# Chat history clearing
# -------------------
def clear_chat_history(user: Optional[str] = None, db_session: Optional[db.SessionLocal] = None):
    """
    Clear chat history rows. If user is None, clears all rows (admin).
    """
    session = db_session or db.SessionLocal()
    try:
        q = session.query(ChatHistory)
        if user:
            q = q.filter(ChatHistory.enteredby == user)
        deleted = q.delete()
        session.commit()
        logger.info("Cleared %s chat entries for user=%s", deleted, user)
        return deleted
    except Exception:
        logger.exception("Failed to clear chat history")
        raise
    finally:
        if db_session is None:
            session.close()
